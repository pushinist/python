from docx import Document

doc = Document()

doc.add_heading("Заголовок документа", 0)

p = doc.add_paragraph("Абзац обычного текста этого документа, ")

r = p.add_run("часть жирным шрифтом").bold = True


r = p.add_run(" а часть курсивом ").italic = True

doc.add_heading('Заголовок первого уровня', level=1)
doc.add_paragraph('Некоторая цитата', style="Intense Quote")
doc.add_paragraph("Элемент ненумерованного списка",
                  style="List Bullet")
doc.add_paragraph("Элемент ненумерованного списка",
                  style='List Number')

table = doc.add_table(rows=1, cols=3)
hdr_cells = table.rows[0].cells
hdr_cells[0].text = 'Номер'
hdr_cells[1].text = 'Название'
hdr_cells[2].text = 'Количество'


doc.save("test.docx")


doc = Document('test.docx')
for par in doc.paragraphs:
    print(par.text)

for table in doc.tables:
    print(table.rows[0].cells[0].text)



from docxtpl import DocxTemplate
import datetime as dt

doc = DocxTemplate('tpl.docx')
context = {
    'name': 'Иван Иванович',
    'event_name': 'танцевально-увеселительное мероприятие .. Кабо-Верде',
    'place': "любой точке Кабо-Верде",
    'date': dt.date.today(),
    'time': dt.datetime.now().strftime("%H:%M"),
    'items': [
        'картинку',
        'корзинку',
        'картонку',
        'маленькую собачонку'
    ]
}
doc.render(context)
doc.save("res.docx")


from pptx import Presentation

prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Тестовый заголовок"
slide.placeholders[1].text = "Тестовый текст"

slide = prs.slides.add_slide(prs.slide_layouts[8])
slide.shapes.title.text = "А теперь с картинкой"
slide.placeholders[1].insert_picture("jpg.jpg")

prs.save("test.pptx")


import xlsxwriter

workbook = xlsxwriter.Workbook('Суммы.xlsx')
worksheet = workbook.add_worksheet()

data = [('Развлечения', 6800), ('Продукты', 25670), ('Транспорт', 3450)]

for row, (item, price) in enumerate(data):
    worksheet.write(row, 0, item)
    worksheet.write(row, 1, price)

row += 1
worksheet.write(row, 0, 'Всего')
worksheet.write(row, 1, '=SUM(B1:B3)')
workbook.close()

workbook = xlsxwriter.Workbook("диаграммы.xlsx")
worksheet = workbook.add_worksheet()

data = [10, 40, 50, 20, 10, 50, 20]
worksheet.write_column('A1', data)

chart = workbook.add_chart({'type': 'pie'})

chart.add_series({'values': '=Sheet1!A1:A7',
                  "categories": "Sheet1!B1:B7"})
worksheet.insert_chart("C1", chart)
workbook.close()

from openpyxl import Workbook, load_workbook
 
wb = Workbook()

ws = wb.active

ws['A1'] = 42

ws.append([1, 2, 3])

import datetime
ws['A2'] = datetime.datetime.now()
wb.save('sample.xlsx')


wb = load_workbook('bred.xlsx')
ws = wb.active

for row in ws.values:
    print(row)

# Задание: взять готовую книгу, создать в ворде заготовку, используя docxtpl и openpyxl. Написать программу, которая создаст документ с нуля, используя pydocx (то же самое)